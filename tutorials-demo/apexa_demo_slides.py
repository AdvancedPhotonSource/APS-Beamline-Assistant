#!/usr/bin/env python3
"""Generate APEXA demo PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Argonne / APS brand-ish colors
DARK_BLUE  = RGBColor(0x00, 0x2B, 0x5C)  # navy
MED_BLUE   = RGBColor(0x00, 0x5E, 0xA2)  # APS blue
ACCENT     = RGBColor(0xE8, 0x7C, 0x00)  # orange accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
RED        = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── helpers ──────────────────────────────────────────────────────────────

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=DARK_GRAY,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_slide_content(tf, items, font_size=16, color=DARK_GRAY, bold_prefix=True):
    """Add bullet points to an existing text frame."""
    for item in items:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        # If item has a bold prefix like "Key: value"
        if bold_prefix and ":" in item and not item.startswith("--"):
            prefix, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = prefix + ":"
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"

def section_header(slide, title, subtitle=""):
    add_bg(slide, DARK_BLUE)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                 title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
                     subtitle, font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

def slide_title_bar(slide, title):
    """Add a colored title bar at the top of a content slide."""
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MED_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8),
                 title, font_size=30, bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "APEXA", font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Advanced Photon EXperiment Assistant", font_size=28, color=RGBColor(0xBB, 0xDE, 0xFB),
             alignment=PP_ALIGN.CENTER)

# Orange accent line
add_shape_bg(slide, Inches(4.5), Inches(3.6), Inches(4.3), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.7),
             "AI-Powered Agentic Framework for Autonomous\nSynchrotron HEDM Experimentation",
             font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Pawan Tripathi  |  Advanced Photon Source  |  Argonne National Laboratory",
             font_size=18, color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "March 2026",
             font_size=16, color=RGBColor(0x64, 0xB5, 0xF6), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "The Problem: Beamtime is Precious")

items = [
    "HEDM data analysis requires memorizing dozens of scripts, flags, and file conventions",
    "MIDAS alone has 20+ executables with complex parameter files (Parameters.txt, 100+ fields)",
    "Operators spend significant time on command-line plumbing instead of science",
    "Motor control requires knowing EPICS PV naming conventions and safety procedures",
    "New users face a steep learning curve \u2014 weeks to become productive",
    "Errors in parameters or workflow order waste precious beamtime",
]

tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5),
                  "", font_size=20, color=DARK_GRAY)
for item in items:
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.space_after  = Pt(10)
    run = p.add_run()
    run.text = "\u2022  " + item
    run.font.size = Pt(20)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is APEXA
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "APEXA: Talk to Your Beamline")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0),
             "You type plain English. APEXA figures out which tools to call, in what order,\n"
             "with what parameters \u2014 and executes them live on your data.",
             font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Example boxes
examples = [
    ("\u201cCalibrate the CeO2 data at 61.332 keV\u201d",
     "Auto-converts energy \u2192 wavelength, finds files,\nruns AutoCalibrateZarr.py, reports beam center & tilts"),
    ("\u201cIntegrate and show me the lineout\u201d",
     "Runs 2D\u21921D integration, auto-detects calibration,\nlaunches MIDAS viewer with correct script & flags"),
    ("\u201cMove sample X to 25.3 mm\u201d",
     "Checks soft limits, verifies limit switches,\nissues move, waits for completion, reports final position"),
]

for i, (query, desc) in enumerate(examples):
    left = Inches(0.5 + i * 4.2)
    # Query box (orange-tinted)
    box = add_shape_bg(slide, left, Inches(3.0), Inches(3.8), Inches(1.2), RGBColor(0xFF, 0xF3, 0xE0))
    add_text_box(slide, left + Inches(0.2), Inches(3.1), Inches(3.4), Inches(1.0),
                 query, font_size=16, bold=True, color=DARK_BLUE)
    # Arrow
    add_text_box(slide, left + Inches(1.5), Inches(4.2), Inches(1), Inches(0.5),
                 "\u2193", font_size=24, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    # Result box (blue-tinted)
    box2 = add_shape_bg(slide, left, Inches(4.6), Inches(3.8), Inches(1.5), RGBColor(0xE3, 0xF2, 0xFD))
    add_text_box(slide, left + Inches(0.2), Inches(4.7), Inches(3.4), Inches(1.3),
                 desc, font_size=14, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Architecture: Multi-Agent Orchestration")

# User query box
add_shape_bg(slide, Inches(4.8), Inches(1.4), Inches(3.7), Inches(0.7), RGBColor(0xE8, 0xEA, 0xF6))
add_text_box(slide, Inches(5.0), Inches(1.45), Inches(3.3), Inches(0.6),
             'User Query (natural language)', font_size=16, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Arrow down
add_text_box(slide, Inches(6.2), Inches(2.1), Inches(1), Inches(0.5),
             "\u2193", font_size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Orchestrator
add_shape_bg(slide, Inches(4.2), Inches(2.5), Inches(5.0), Inches(0.8), ACCENT)
add_text_box(slide, Inches(4.4), Inches(2.55), Inches(4.6), Inches(0.7),
             'Orchestrator Agent (keyword-score routing)', font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Arrow down
add_text_box(slide, Inches(6.2), Inches(3.3), Inches(1), Inches(0.4),
             "\u2193", font_size=28, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

# 5 specialist agents
agents = [
    ("Calibration\nAgent", "T=0.3", RGBColor(0xE8, 0xF5, 0xE9)),
    ("Analysis\nAgent", "T=0.5", RGBColor(0xE3, 0xF2, 0xFD)),
    ("Knowledge\nAgent", "T=0.6", RGBColor(0xFD, 0xF3, 0xE7)),
    ("Visualization\nAgent", "T=0.3", RGBColor(0xF3, 0xE5, 0xF5)),
    ("Motor\nAgent", "T=0.2", RGBColor(0xFF, 0xEB, 0xEE)),
]

for i, (name, temp, color) in enumerate(agents):
    left = Inches(0.4 + i * 2.6)
    add_shape_bg(slide, left, Inches(3.8), Inches(2.3), Inches(1.2), color)
    add_text_box(slide, left + Inches(0.1), Inches(3.85), Inches(2.1), Inches(0.8),
                 name, font_size=15, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(4.6), Inches(2.1), Inches(0.3),
                 temp, font_size=11, color=RGBColor(0x75, 0x75, 0x75), alignment=PP_ALIGN.CENTER)

# Arrows down from agents
for i in range(5):
    add_text_box(slide, Inches(1.2 + i * 2.6), Inches(5.0), Inches(0.5), Inches(0.4),
                 "\u2193", font_size=22, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

# MCP Tool Servers
add_shape_bg(slide, Inches(0.4), Inches(5.4), Inches(3.8), Inches(0.8), MED_BLUE)
add_text_box(slide, Inches(0.6), Inches(5.45), Inches(3.4), Inches(0.7),
             'Core Server (9 tools)\nFiles, Shell, X-ray Calc', font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4.6), Inches(5.4), Inches(4.2), Inches(0.8), MED_BLUE)
add_text_box(slide, Inches(4.8), Inches(5.45), Inches(3.8), Inches(0.7),
             'MIDAS Server (21 tools)\nCalibration, Integration, FF/NF/PF-HEDM', font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(9.2), Inches(5.4), Inches(3.7), Inches(0.8), MED_BLUE)
add_text_box(slide, Inches(9.4), Inches(5.45), Inches(3.3), Inches(0.7),
             'EPICS Motor Server (12 tools)\nMotor Control via Channel Access', font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# Bottom: Argo Gateway
add_shape_bg(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.7), DARK_BLUE)
add_text_box(slide, Inches(2.7), Inches(6.52), Inches(7.9), Inches(0.6),
             'Argo Gateway API  \u2014  GPT-4o / GPT-5 / Claude Opus / Gemini 2.5  (switchable at runtime)',
             font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — 42 Live Tools
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "42 Live Tools Across 3 MCP Servers")

# Core Server column
add_shape_bg(slide, Inches(0.4), Inches(1.5), Inches(3.8), Inches(0.6), MED_BLUE)
add_text_box(slide, Inches(0.6), Inches(1.55), Inches(3.4), Inches(0.5),
             'Core Server (9 tools)', font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

core_tools = [
    "list_directory \u2014 Browse experiment folders",
    "read_file / write_file \u2014 File I/O",
    "run_command \u2014 Shell execution",
    "xray_calculate \u2014 Bragg's law, d-spacing, energy\u2194wavelength",
    "validate_beamline_parameters",
    "list_common_calibrants",
    "check_environment",
]
tf = add_text_box(slide, Inches(0.5), Inches(2.2), Inches(3.7), Inches(4.5),
                  "", font_size=13, color=DARK_GRAY)
for t in core_tools:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = "\u2022 " + t
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"

# MIDAS Server column
add_shape_bg(slide, Inches(4.6), Inches(1.5), Inches(4.3), Inches(0.6), MED_BLUE)
add_text_box(slide, Inches(4.8), Inches(1.55), Inches(3.9), Inches(0.5),
             'MIDAS Server (21 tools)', font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

midas_tools = [
    "midas_auto_calibrate",
    "midas_integrate_2d_to_1d",
    "midas_batch_integrate",
    "run_ff_hedm_full_workflow",
    "run_nf_hedm_reconstruction",
    "run_pf_hedm_workflow",
    "match_grains (Hungarian alg.)",
    "calculate_misorientation",
    "run_forward_simulation",
    "extract_grain_centroids",
    "convert_nf_to_dream3d",
    "overlay_ff_nf_results",
    "create_midas_parameter_file",
    "get_material_properties",
    "estimate_parameters_from_image",
]
tf = add_text_box(slide, Inches(4.7), Inches(2.2), Inches(4.1), Inches(5.0),
                  "", font_size=13, color=DARK_GRAY)
for t in midas_tools:
    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = "\u2022 " + t
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"

# Motor Server column
add_shape_bg(slide, Inches(9.3), Inches(1.5), Inches(3.7), Inches(0.6), MED_BLUE)
add_text_box(slide, Inches(9.5), Inches(1.55), Inches(3.3), Inches(0.5),
             'Motor Server (12 tools)', font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

motor_tools = [
    "get_motor_position",
    "get_motor_status",
    "move_motor_absolute",
    "move_motor_relative",
    "stop_motor",
    "set_motor_velocity",
    "jog_motor",
    "tweak_motor",
    "get_motor_limits",
    "set_motor_limits",
    "list_motors",
    "home_motor",
]
tf = add_text_box(slide, Inches(9.4), Inches(2.2), Inches(3.5), Inches(4.5),
                  "", font_size=13, color=DARK_GRAY)
for t in motor_tools:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = "\u2022 " + t
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Motor Control & Safety
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Motor Control: AI with Built-in Safety")

# Left side: capabilities
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "Natural-Language Motor Control via EPICS", font_size=22, bold=True, color=DARK_BLUE)

cap_items = [
    '"What\'s the position of motor m1?" \u2192 reads RBV, VAL, units',
    '"Move m1 to 50.0" \u2192 checks limits, issues move, waits, reports',
    '"Jog forward for 2 seconds" \u2192 timed jog with auto-stop',
    '"Tweak m3 forward by 0.01" \u2192 fine step via TWV/TWF fields',
    '"Stop!" \u2192 immediate STOP=1, always allowed, no questions asked',
    "Works with ANY IOC using standard motorRecord convention",
]
tf = add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(3.5),
                  "", font_size=16, color=DARK_GRAY)
for item in cap_items:
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = "\u2022  " + item
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"

# Right side: safety box
add_shape_bg(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.0), RGBColor(0xFF, 0xEB, 0xEE))

add_text_box(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5),
             "Safety Policy (never bypassed)", font_size=20, bold=True, color=RED)

safety_items = [
    "Soft-limit check (HLM/LLM) before every move \u2014 rejects out-of-range targets",
    "Hard limit-switch guard (HLS/LLS) \u2014 refuses motion when at limit",
    "Large-move protection \u2014 moves >50% of travel range require explicit confirmation",
    "Never sets STOP=0 (arming is not the AI's job)",
    "Never homes a motor without explicit user instruction",
    "Motor Agent temperature = 0.2 (most deterministic of all agents)",
]
tf = add_text_box(slide, Inches(7.5), Inches(2.3), Inches(5.0), Inches(4.0),
                  "", font_size=15, color=DARK_GRAY)
for item in safety_items:
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = "\u2022  " + item
    run.font.size = Pt(15)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Agent Skills
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Agent Skills: Codified Beamline Knowledge")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.7),
             "Structured knowledge documents that guide each specialist agent \u2014\n"
             "version-locked to the MIDAS manual so the AI stays accurate when MIDAS updates.",
             font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

skills = [
    ("midas-calibrate", "Exact v10 calibration flags, output filenames,\nerror patterns, convergence criteria"),
    ("midas-integrate", "Updated to v11.0 manual, 2D\u21921D workflow,\nauto-detect calibration files"),
    ("midas-hedm", "FF/NF/PF-HEDM full workflow references,\nGPU executables, checkpoint/resume"),
    ("midas-visualize", "Plot types, file format \u2192 viewer script mapping,\ncritical flags (case-sensitive!)"),
]

for i, (name, desc) in enumerate(skills):
    top = Inches(2.6 + i * 1.15)
    add_shape_bg(slide, Inches(0.8), top, Inches(3.0), Inches(0.9), RGBColor(0xE8, 0xF5, 0xE9))
    add_text_box(slide, Inches(1.0), top + Inches(0.15), Inches(2.6), Inches(0.6),
                 name, font_size=20, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.2), top + Inches(0.1), Inches(8.5), Inches(0.8),
                 desc, font_size=16, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Smart Features
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Smart Features for Beamtime Efficiency")

features = [
    ("Proactive Suggestions", "After each analysis step, APEXA suggests logical next steps\n(e.g., after integration \u2192 \u201cidentify phases, check for peak splitting\u201d)"),
    ("Real-Time Monitoring", "Watch directories for new images, auto-analyze quality,\nalert on saturation, low SNR, or missing diffraction rings"),
    ("Image Quality Analysis", "Automated SNR, saturation detection, hot pixel detection,\nring detection \u2014 flags problems before you waste time analyzing bad data"),
    ("Session Persistence", "Save/load experiment sessions across restarts \u2014 tracks\nanalysis history, key findings, active files"),
    ("Smart Caching", "Expensive read-only operations cached to avoid\nredundant calls \u2014 faster repeat queries"),
    ("Error Prevention", "Validates directories, parameter files, and inputs before\nexecution \u2014 catches mistakes before they waste beamtime"),
]

for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.5 + row * 1.8)

    add_shape_bg(slide, left, top, Inches(6.0), Inches(1.5), LIGHT_GRAY)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), Inches(5.6), Inches(0.4),
                 title, font_size=18, bold=True, color=MED_BLUE)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55), Inches(5.6), Inches(0.9),
                 desc, font_size=14, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Multi-Model Support
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Multi-Model Support via Argo Gateway")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
             "Swap AI models at runtime without changing workflows \u2014 all through Argonne\u2019s Argo API",
             font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

models = [
    ("OpenAI", "GPT-4o, GPT-4 Turbo, GPT-5, GPT-5 Mini/Nano", "PROD / DEV"),
    ("Anthropic", "Claude Opus 4/4.1, Sonnet 4/4.5/3.7", "DEV"),
    ("Google", "Gemini 2.5 Pro, Gemini 2.5 Flash", "DEV"),
]

for i, (provider, model_list, env) in enumerate(models):
    top = Inches(2.5 + i * 1.2)
    color = [RGBColor(0xE3, 0xF2, 0xFD), RGBColor(0xFD, 0xF3, 0xE7), RGBColor(0xE8, 0xF5, 0xE9)][i]
    add_shape_bg(slide, Inches(1.5), top, Inches(10.3), Inches(0.9), color)
    add_text_box(slide, Inches(1.7), top + Inches(0.15), Inches(2.2), Inches(0.6),
                 provider, font_size=20, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(4.0), top + Inches(0.15), Inches(5.5), Inches(0.6),
                 model_list, font_size=16, color=DARK_GRAY)
    add_text_box(slide, Inches(9.8), top + Inches(0.15), Inches(1.8), Inches(0.6),
                 env, font_size=14, bold=True, color=RGBColor(0x75, 0x75, 0x75), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2),
             "3 Interfaces: CLI (start_beamline_assistant.sh), Gradio Web UI, Web Server API\n"
             "All share the same run_query() entry point \u2014 identical behavior across interfaces",
             font_size=17, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Demo Scenarios
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Live Demo Scenarios")

demos = [
    ("1. Calibration", "\"Calibrate CeO2 data in test5 at 61.332 keV\"\n\u2192 energy\u2192wavelength \u2192 find files \u2192 AutoCalibrateZarr \u2192 report BC, Lsd, tilts"),
    ("2. Integration + Viz", "\"Integrate the CeO2 data\" then \"Show the lineout\"\n\u2192 2D\u21921D integration \u2192 auto-launch MIDAS viewer with correct script"),
    ("3. X-ray Calculations", "\"Calculate d-spacing for (110) in bcc iron\"\n\u2192 xray_calculate tool \u2014 never computes manually"),
    ("4. Motor Control", "\"Move m1 to 50.0\" \u2192 limit check \u2192 move \u2192 wait \u2192 report\n\"Stop!\" \u2192 immediate STOP=1"),
    ("5. Full HEDM Workflow", "\"Run FF-HEDM on test5\"\n\u2192 validate dir \u2192 run_ff_hedm_full_workflow \u2192 grains, convergence, outputs"),
    ("6. Multi-turn", "Remembers context: \"Now track grains between step 1 and 2\"\n\u2192 knows what data you were working with"),
]

for i, (title, desc) in enumerate(demos):
    row = i // 2
    col = i % 2
    left = Inches(0.4 + col * 6.5)
    top = Inches(1.5 + row * 1.85)

    add_shape_bg(slide, left, top, Inches(6.1), Inches(1.6), LIGHT_GRAY)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.08), Inches(5.7), Inches(0.4),
                 title, font_size=18, bold=True, color=MED_BLUE)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), Inches(5.7), Inches(1.0),
                 desc, font_size=14, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Impact / Why It Matters
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Why It Matters for Beamline Users")

impacts = [
    ("No more memorizing flags", "You say \"calibrate\" \u2014 APEXA knows AutoCalibrateZarr.py --paramFN ... --darkFN ..."),
    ("No more wrong script paths", "Agent Skills are version-locked to the MIDAS v10/v11 manual"),
    ("Errors caught before execution", "Validates directories, parameters, and limits before running anything"),
    ("Faster beamtime", "Operators spend time on science, not command-line plumbing"),
    ("Safe motor control", "AI respects physical constraints; humans confirm large moves"),
    ("Works with your model of choice", "Swap between GPT-5, Claude, Gemini without changing workflows"),
]

for i, (title, desc) in enumerate(impacts):
    top = Inches(1.5 + i * 0.9)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), top + Inches(0.05), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MED_BLUE
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    run_c = tf_c.paragraphs[0].add_run()
    run_c.text = str(i + 1)
    run_c.font.size = Pt(16)
    run_c.font.bold = True
    run_c.font.color.rgb = WHITE

    add_text_box(slide, Inches(1.5), top, Inches(3.5), Inches(0.5),
                 title, font_size=18, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(5.2), top + Inches(0.02), Inches(7.5), Inches(0.5),
                 desc, font_size=16, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Command Reference: Calibration & Integration
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Command Reference: Calibration & Integration")

# Left column — Calibration
add_shape_bg(slide, Inches(0.3), Inches(1.4), Inches(6.3), Inches(0.5), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.5), Inches(1.42), Inches(6.0), Inches(0.45),
             "Calibration", font_size=18, bold=True, color=DARK_BLUE)

cal_cmds = [
    'Calibrate the CeO2 image in test5',
    'Calibrate the CeO2 data at 61.332 keV in test5',
    'Calibrate test5/CeO_000001.tif.ge using\n  test5/Parameters.txt with dark file\n  test5/dark_CeO_000001.tif.ge',
    'What calibrants are available?',
    'Validate: energy 61.332 keV, detector\n  distance 650 mm, beam center 1024 1024',
]
tf = add_text_box(slide, Inches(0.4), Inches(1.95), Inches(6.1), Inches(4.5),
                  "", font_size=12, color=DARK_GRAY)
for cmd in cal_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"

# Right column — Integration
add_shape_bg(slide, Inches(6.9), Inches(1.4), Inches(6.1), Inches(0.5), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(7.1), Inches(1.42), Inches(5.8), Inches(0.45),
             "Integration (2D \u2192 1D)", font_size=18, bold=True, color=DARK_BLUE)

int_cmds = [
    'Integrate the diffraction image in test5',
    'Integrate test5/CeO_000001.tif.ge using\n  calibration file refined_MIDAS_params_CeO.txt',
    'Integrate test5/CeO_000001.tif.ge with\n  dark file test5/dark_CeO_000001.tif.ge',
    'Batch integrate data/sample_003083.ge1.h5\n  frames 3083 to 3085 using dark file\n  data/dark_003084.ge1.h5 with 80 CPUs',
]
tf = add_text_box(slide, Inches(7.0), Inches(1.95), Inches(6.0), Inches(4.5),
                  "", font_size=12, color=DARK_GRAY)
for cmd in int_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"

# Bottom — X-ray Calculations
add_shape_bg(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(0.5), RGBColor(0xFD, 0xF3, 0xE7))
add_text_box(slide, Inches(0.5), Inches(5.02), Inches(5.0), Inches(0.45),
             "X-ray Calculations", font_size=18, bold=True, color=DARK_BLUE)

xray_cmds = [
    ('Convert 61.332 keV to wavelength', 'Calculate d-spacing for (110) plane in bcc iron'),
    ('Convert 0.2022 angstroms to energy', 'What is 2-theta for d=2.03 \u00c5 at 61.332 keV?'),
    ('Calculate strain: measured d 2.035, ref 2.028', 'List all available materials'),
]
for i, (left_cmd, right_cmd) in enumerate(xray_cmds):
    y = Inches(5.55 + i * 0.35)
    tf_l = add_text_box(slide, Inches(0.4), y, Inches(6.1), Inches(0.3),
                        "", font_size=11, color=DARK_GRAY)
    p = tf_l.paragraphs[0]
    run = p.add_run()
    run.text = "\u25b6  " + left_cmd
    run.font.size = Pt(11)
    run.font.name = "Consolas"
    run.font.color.rgb = DARK_GRAY

    tf_r = add_text_box(slide, Inches(7.0), y, Inches(6.0), Inches(0.3),
                        "", font_size=11, color=DARK_GRAY)
    p = tf_r.paragraphs[0]
    run = p.add_run()
    run.text = "\u25b6  " + right_cmd
    run.font.size = Pt(11)
    run.font.name = "Consolas"
    run.font.color.rgb = DARK_GRAY


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Command Reference: HEDM Workflows & Grain Analysis
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Command Reference: HEDM & Grain Analysis")

# Left — HEDM Workflows
add_shape_bg(slide, Inches(0.3), Inches(1.4), Inches(6.3), Inches(0.5), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(0.5), Inches(1.42), Inches(6.0), Inches(0.45),
             "HEDM Workflows", font_size=18, bold=True, color=DARK_BLUE)

hedm_cmds = [
    'Run FF-HEDM reconstruction on test5',
    'Run FF-HEDM on test5 with GPU enabled',
    'Resume FF-HEDM from checkpoint test5/checkpoint.h5',
    'Restart FF-HEDM from the indexing step in test5',
    'Run NF-HEDM reconstruction using\n  test5/Parameters.txt',
    'Run NF-HEDM using test5/Parameters.txt with\n  FF grains from test5/Grains.csv',
    'Run PF-HEDM workflow using Parameters.txt\n  with positions file scan_positions.csv',
]
tf = add_text_box(slide, Inches(0.4), Inches(1.95), Inches(6.1), Inches(5.0),
                  "", font_size=12, color=DARK_GRAY)
for cmd in hedm_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"

# Right — Grain Analysis
add_shape_bg(slide, Inches(6.9), Inches(1.4), Inches(6.1), Inches(0.5), RGBColor(0xF3, 0xE5, 0xF5))
add_text_box(slide, Inches(7.1), Inches(1.42), Inches(5.8), Inches(0.45),
             "Grain Analysis & Post-Processing", font_size=18, bold=True, color=DARK_BLUE)

grain_cmds = [
    'Match grains between step1/Grains.csv and\n  step2/Grains.csv with position tolerance\n  100 microns and orientation tolerance 2 deg',
    'Calculate misorientation between grain 1 and\n  grain 5 in test5/Grains.csv for FCC',
    'Extract grain centroids from test5/Grains.mic\n  with minimum grain size 100 voxels',
    'Run forward simulation using test5/Grains.csv\n  and test5/Parameters.txt',
    'Convert NF-HEDM results to Dream3D format',
    'Overlay FF and NF results from test5',
]
tf = add_text_box(slide, Inches(7.0), Inches(1.95), Inches(6.0), Inches(5.0),
                  "", font_size=12, color=DARK_GRAY)
for cmd in grain_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Command Reference: Visualization & Motor Control
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Command Reference: Visualization & Motor Control")

# Left — Visualization
add_shape_bg(slide, Inches(0.3), Inches(1.4), Inches(6.3), Inches(0.5), RGBColor(0xF3, 0xE5, 0xF5))
add_text_box(slide, Inches(0.5), Inches(1.42), Inches(6.0), Inches(0.45),
             "Visualization", font_size=18, bold=True, color=DARK_BLUE)

viz_cmds = [
    'Show me the diffraction image in test5',
    'Plot the lineout for test5 integration',
    'Show the caked output in test5/integration',
    'Plot the calibration results in test5',
    'Launch the live viewer for\n  test5/CeO_lineout.bin with 2000 radial bins',
    'Show the FF-HEDM grain results in test5',
    'View the NF microstructure map test5/Grains.mic',
    'Compare lineouts from test5 and test6',
]
tf = add_text_box(slide, Inches(0.4), Inches(1.95), Inches(6.1), Inches(5.0),
                  "", font_size=12, color=DARK_GRAY)
for cmd in viz_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"

# Right — Motor Control
add_shape_bg(slide, Inches(6.9), Inches(1.4), Inches(6.1), Inches(0.5), RGBColor(0xFF, 0xEB, 0xEE))
add_text_box(slide, Inches(7.1), Inches(1.42), Inches(5.8), Inches(0.45),
             "Motor Control (EPICS)", font_size=18, bold=True, color=DARK_BLUE)

motor_cmds = [
    'What is the position of motor m1?',
    'Show me the status of motor m1',
    'Show positions of motors m1, m2, m3, m4',
    'Move m1 to 25.3',
    'Move m1 by +0.5',
    'Tweak m1 forward by 0.01',
    'Jog m1 forward for 2 seconds',
    'Stop motor m1!',
    'Set motor m1 velocity to 5.0',
    'What are the limits of m1?',
    'Set high limit of m1 to 200 and low to -50',
    'Home motor m1 in the forward direction',
]
tf = add_text_box(slide, Inches(7.0), Inches(1.95), Inches(6.0), Inches(5.0),
                  "", font_size=12, color=DARK_GRAY)
for cmd in motor_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Command Reference: Knowledge & Multi-Turn
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title_bar(slide, "Command Reference: Knowledge & Multi-Turn Workflows")

# Left — Knowledge
add_shape_bg(slide, Inches(0.3), Inches(1.4), Inches(6.3), Inches(0.5), RGBColor(0xFD, 0xF3, 0xE7))
add_text_box(slide, Inches(0.5), Inches(1.42), Inches(6.0), Inches(0.45),
             "Knowledge & Material Properties", font_size=18, bold=True, color=DARK_BLUE)

know_cmds = [
    'What are the material properties of CeO2?',
    'Give me the lattice parameters for LaB6',
    'What is the space group of titanium?',
    'Material properties of Steel_316L',
    'What is the difference between FF and NF HEDM?',
    'What are best practices for calibration\n  at high energy?',
    'What are typical HEDM parameters for 61 keV?',
    'Estimate detector distance from ring radii\n  [412, 478, 675] px at 0.2022 \u00c5, pixel 200 \u00b5m',
]
tf = add_text_box(slide, Inches(0.4), Inches(1.95), Inches(6.1), Inches(5.0),
                  "", font_size=12, color=DARK_GRAY)
for cmd in know_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"

# Right — Multi-Turn Workflow
add_shape_bg(slide, Inches(6.9), Inches(1.4), Inches(6.1), Inches(0.5), RGBColor(0xE8, 0xEA, 0xF6))
add_text_box(slide, Inches(7.1), Inches(1.42), Inches(5.8), Inches(0.45),
             "Multi-Turn Workflow (context-aware)", font_size=18, bold=True, color=DARK_BLUE)

multi_cmds = [
    '1. Calibrate the CeO2 data in test5\n     at 61.332 keV',
    '2. Now integrate that data',
    '3. Show me the lineout',
    '4. What phases match these peaks?',
    '5. Run FF-HEDM on this data',
    '6. How many grains were found?',
    '7. Track grains between step 1 and step 2',
    '8. Export the results to Dream3D',
]
tf = add_text_box(slide, Inches(7.0), Inches(1.95), Inches(6.0), Inches(5.0),
                  "", font_size=12, color=DARK_GRAY)
for cmd in multi_cmds:
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = "\u25b6  " + cmd
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_GRAY
    run.font.name = "Consolas"

# Bottom note
add_text_box(slide, Inches(6.9), Inches(6.2), Inches(6.1), Inches(0.8),
             "APEXA remembers context across turns \u2014\n"
             "no need to repeat file paths or parameters.",
             font_size=14, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — Closing
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
             "APEXA", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(5.0), Inches(2.9), Inches(3.3), Inches(0.05), ACCENT)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "Talk to your beamline.", font_size=28, color=RGBColor(0xBB, 0xDE, 0xFB),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(1.0),
             "42 tools  |  5 specialist agents  |  3 MCP servers\n"
             "Calibration \u2022 Integration \u2022 FF/NF/PF-HEDM \u2022 Motor Control \u2022 Visualization",
             font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Pawan Tripathi  |  ptripathi@anl.gov  |  Advanced Photon Source, Argonne National Laboratory",
             font_size=18, color=RGBColor(0x90, 0xCA, 0xF9), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "github.com/AdvancedPhotonSource/APS-Beamline-Assistant",
             font_size=16, color=RGBColor(0x64, 0xB5, 0xF6), alignment=PP_ALIGN.CENTER)

# ── Save ──
output = "tutorials-demo/APEXA_Demo_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
