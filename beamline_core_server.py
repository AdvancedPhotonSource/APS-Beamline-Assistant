#!/usr/bin/env python3
"""
Beamline Core Operations MCP Server
Unified server combining filesystem, command execution, and X-ray utilities

This server consolidates:
- filesystem_server.py → File operations
- command_executor_server.py → Safe command execution
- analysis_utilities_server.py → X-ray calculations and utilities

Author: Beamline Assistant Team
Organization: Argonne National Laboratory
Created: 2025 (Unified architecture)
"""

from typing import Any, Optional, List, Dict
import json
import sys
import os
import logging
from pathlib import Path

# Suppress noisy third-party library startup messages
logging.getLogger("numexpr").setLevel(logging.WARNING)
logging.getLogger("numexpr.utils").setLevel(logging.WARNING)
import shlex
import subprocess
import stat
import time
import platform
import shutil
import numpy as np
import logging
from mcp.server.fastmcp import FastMCP

from apexa_remote_exec import remote_run, ssh_hint

# Suppress verbose MCP server logging
logging.getLogger("mcp").setLevel(logging.WARNING)
logging.getLogger("fastmcp").setLevel(logging.WARNING)

# =============================================================================
# INITIALIZATION
# =============================================================================

mcp = FastMCP("beamline-core")

# Try to import X-ray utilities (domain-specific package)
try:
    import xrayutilities as xu
    import xrayutilities.materials as xumaterials
    XRAYUTILITIES_AVAILABLE = True
except ImportError:
    XRAYUTILITIES_AVAILABLE = False
    print("⚠️ xrayutilities not available - install with: pip install xrayutilities", file=sys.stderr)

# Try to import optional dependencies for image analysis
try:
    import fabio
    from scipy.signal import find_peaks
    SCIPY_AVAILABLE = True
except ImportError:
    SCIPY_AVAILABLE = False
    print("⚠️ scipy/fabio not available - some image analysis features limited", file=sys.stderr)

# =============================================================================
# SECURITY: ALLOWED COMMANDS
# =============================================================================

ALLOWED_COMMANDS = {
    # Shell interpreters — enables bash -c "..." multi-command inline scripts,
    # making run_command functionally equivalent to Claude Code's Bash tool.
    # Real safety guarantee is the motor server's hardware gate, not this list.
    'bash', 'sh', 'zsh',
    # Basic system commands
    'ls', 'pwd', 'cat', 'head', 'tail', 'less', 'more',
    'grep', 'find', 'wc', 'sort', 'uniq', 'diff', 'sed', 'awk',
    'cp', 'mv', 'mkdir', 'rmdir', 'touch', 'chmod', 'ln', 'readlink',
    # Deletion (rm/unlink; rmdir above). ENABLED so an APPROVED deletion actually
    # executes — the real safeguard is the client-side human permission gate in
    # APEXAClient.execute_tool_call (this server subprocess cannot prompt a human),
    # not this allowlist. Same spirit as the motor server's hardware gate.
    'rm', 'unlink',
    'echo', 'date', 'hostname', 'uname', 'whoami', 'env', 'which',
    'file', 'stat', 'tree', 'du', 'df', 'free', 'ps', 'nproc', 'lscpu', 'uptime',
    # Read-only text/binary inspection + comparison utilities (safe, no mutation)
    'md5sum', 'sha1sum', 'sha256sum', 'cksum', 'cmp', 'comm', 'cut', 'paste',
    'tr', 'tee', 'xargs', 'column', 'basename', 'dirname', 'realpath', 'seq',
    'tac', 'nl', 'od', 'xxd', 'hexdump', 'jq',
    'tar', 'gzip', 'gunzip', 'zip', 'unzip',
    'curl', 'wget', 'ssh', 'scp', 'rsync',
    'python', 'python3', 'pip', 'pip3', 'uv', 'conda', 'git',
    'make', 'cmake', 'gcc', 'g++',
    'caget', 'caput', 'camonitor', 'cainfo',

    # ── FF-HEDM CPU executables ──────────────────────────────────────────────
    'GetHKLListZarr',
    'PeaksFittingOMPZarrRefactor',
    'MergeOverlappingPeaksAllZarr',
    'CalcRadiusAllZarr',
    'FitSetupZarr',
    'SaveBinData',
    'SaveBinDataScanning',
    'IndexerOMP',
    'IndexerScanningOMP',
    'FitPosOrStrainsOMP',
    'FitOrStrainsScanningOMP',
    'FitPosOrStrainsDoubleDataset',
    'FitMultipleGrains',
    'FitWedgeParallel',
    'ProcessGrains',
    'ProcessGrainsScanningHEDM',
    'MatchGrains',
    'MergeMultipleScans',
    'mergeScansScanning',

    # ── FF-HEDM GPU executables (v10) ────────────────────────────────────────
    'IndexerGPU',
    'IndexerScanningGPU',
    'FitPosOrStrainsGPU',
    'FitOrStrainsScanningGPU',
    'IntegratorFitPeaksGPUStream',

    # ── FF-HEDM calibration & integration ───────────────────────────────────
    'CalibrantIntegratorOMP',       # primary calibrant integrator (v10)
    'CalibrantPanelShiftsOMP',      # multi-panel shift calibration (v10)
    'CalibrantOMP',                 # legacy (keep for compatibility)
    'FitTiltBCLsdSample',
    'IntegratorZarrOMP',            # primary integrator (v10)
    'IntegratorZarrOMP_f32',        # float32 variant
    'DetectorMapper',
    'MapDatasets',
    'ForwardSimulationCompressed',
    'findSingleSolutionPFRefactored',
    'findMultipleSolutionsPF',
    'FindSaturatedPixels',
    'CalcRadius',

    # ── NF-HEDM executables ──────────────────────────────────────────────────
    'GetHKLListNF',
    'MakeHexGrid',
    'MakeDiffrSpots',
    'MedianImageLibTiff',
    'ProcessImagesCombined',        # replaces ImageProcessingLibTiffOMP (v10)
    'MMapImageInfo',
    'FitOrientationOMP',
    'FitOrientationGPU',            # GPU NF reconstruction (v10)
    'FitOrientationParameters',
    'FitOrientationParametersMultiPoint',
    'GenSeedOrientationsFF2NFHEDM',
    'SimulateDiffractionSpots',
    'filterGridfromTomo',
    'ParseMic',
    'ParseDeconvOutput',
    'Mic2GrainsList',
    'NFGrainCentroids',
    'compareNF',
    'simulateNF',

    # ── TOMO executables ─────────────────────────────────────────────────────
    'GenMedianDark',

    # ── Python workflow drivers ──────────────────────────────────────────────
    'ff_MIDAS.py',
    'nf_MIDAS.py',
    'nf_MIDAS_Multiple_Resolutions.py',
    'pf_MIDAS.py',
    'integrator.py',
    'integrator_server.py',
    'integrator_batch_process.py',
    'AutoCalibrateZarr.py',
    'ffGenerateZipRefactor.py',
    'ffGenerateZip.py',
    'runDTrecon.py',

    # ── Python utilities (v10) ───────────────────────────────────────────────
    'match_grains.py',
    'phase_id.py',
    'validate_calibration.py',
    'generate_mask.py',
    'undistort_image.py',
    'updateZarrDset.py',
    'extract_lineouts.py',
    'fit_caked_peaks.py',
    'BatchCake.py',
    'runScanning.py',
    'evalScanning.py',
    'calcMiso.py',
    'simulatePeaks.py',
    'blobPeaksearch.py',
    'nf_paraview_gen.py',
    'nf_mic_to_grains.py',
    'NFGrainCentroids.py',
    'DL2FF.py',
    'GE2Tiff.py',
    'hdf_gen_nf.py',
    'PlotFFNF.py',
}

# Case-insensitive lookup set. `_base_cmd` lowercases the executable token, but
# ALLOWED_COMMANDS above stores canonical mixed case (AutoCalibrateZarr.py,
# ff_MIDAS.py, IndexerOMP, ProcessGrains, …). Comparing a lowercased token
# against the mixed-case set NEVER matches, so every mixed-case MIDAS binary was
# wrongly refused when invoked by name (only all-lowercase entries like python3
# passed). Match against this lowercased copy instead. ALLOWED_COMMANDS stays
# the source of truth and is still what `available_commands` displays.
_ALLOWED_COMMANDS_LOWER = {c.lower() for c in ALLOWED_COMMANDS}

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def format_result(result: dict) -> str:
    """Format results as JSON string"""
    return json.dumps(result, indent=2, default=str)

def format_file_info(path: Path) -> dict:
    """Get detailed file information"""
    try:
        stat_info = path.stat()
        return {
            "name": path.name,
            "path": str(path.absolute()),
            "size": stat_info.st_size,
            "modified": time.ctime(stat_info.st_mtime),
            "created": time.ctime(stat_info.st_ctime),
            "is_file": path.is_file(),
            "is_directory": path.is_dir(),
            "permissions": stat.filemode(stat_info.st_mode),
            "owner_readable": bool(stat_info.st_mode & stat.S_IRUSR),
            "owner_writable": bool(stat_info.st_mode & stat.S_IWUSR),
            "owner_executable": bool(stat_info.st_mode & stat.S_IXUSR)
        }
    except Exception as e:
        return {"error": str(e)}

def _base_cmd(token: str) -> str:
    """Return the bare command name from a full path or bare name."""
    return Path(token).name.lower()


# Shell interpreters whose -c argument is validated recursively.
_SHELL_INTERPRETERS = {'bash', 'sh', 'zsh', 'dash', 'ksh'}

# Shell builtins and control-flow keywords that are always allowed inside
# bash -c scripts.  These are not external executables — they are part of
# the shell itself and cannot be blocked by an allowlist anyway.
_SHELL_BUILTINS = {
    'for', 'while', 'until', 'do', 'done', 'if', 'then', 'else', 'elif',
    'fi', 'case', 'esac', 'in', 'function', 'return', 'break', 'continue',
    'exit', 'cd', 'export', 'source', 'local', 'true', 'false', 'test',
    '[', '[[', ']]', 'read', 'declare', 'typeset', 'let', 'set', 'unset',
    'shift', 'eval', 'printf', 'time', 'exec', 'wait', 'jobs', 'kill',
}


def _split_on_unquoted_ops(command: str) -> list:
    """Split a command into pipeline/list segments on shell operators (| || && ; &
    and newlines) that appear OUTSIDE single/double quotes. Content inside quotes
    (e.g. an awk/sed program) is preserved intact, so its operators and $-fields
    are never mistaken for command separators or command names."""
    segments, buf = [], []
    quote = None            # None, "'", or '"'
    i, n = 0, len(command)
    while i < n:
        c = command[i]
        if quote:
            buf.append(c)
            if c == quote:
                quote = None
            i += 1
            continue
        if c in ("'", '"'):
            quote = c; buf.append(c); i += 1; continue
        # Two-char operators
        if c in '|&' and i + 1 < n and command[i + 1] == c:   # || or &&
            segments.append(''.join(buf)); buf = []; i += 2; continue
        if c in '|;&\n':                                        # | ; & newline
            segments.append(''.join(buf)); buf = []; i += 1; continue
        buf.append(c); i += 1
    segments.append(''.join(buf))
    return segments


def _validate_segments(command: str) -> tuple:
    """Split a shell command on operators and validate each executable segment.

    Returns (allowed: bool, blocked: list[str]).
    Handles redirections (>, >>, <, 2>) by stripping them before validation.
    Does NOT handle shell-quoted strings specially — call this only on content
    that has already been extracted from quoted arguments (e.g. the -c arg).
    """
    import re as _re
    blocked = []
    # Strip heredoc bodies (<<MARKER … MARKER) BEFORE splitting. Embedded report
    # text / Python / markdown tables contain |, &, ;, and barewords (e.g. a
    # table row "| 1. Peak | … |") that must NOT be parsed as shell pipeline
    # segments — that was the "Command not allowed: 1" false block.
    def _strip_heredocs(s: str) -> str:
        lines = s.split('\n')
        out, i = [], 0
        while i < len(lines):
            m = _re.search(r'<<-?\s*["\']?([A-Za-z_]\w*)["\']?', lines[i])
            if m:
                out.append(lines[i][:m.start()])      # keep text before <<
                marker = m.group(1)
                i += 1
                while i < len(lines) and lines[i].strip() != marker:
                    i += 1                              # skip heredoc body
                i += 1                                  # skip closing marker
                continue
            out.append(lines[i])
            i += 1
        return '\n'.join(out)

    command = _strip_heredocs(command)
    # Split on shell operators AND newlines so each real command line is checked
    # (catches e.g. a destructive command on its own line, not just the first),
    # but ONLY on operators that are OUTSIDE quotes. A naive regex split breaks
    # quoted program text — e.g. awk '{print $5}' or sed 's/;/,/g' — into bogus
    # "segments" whose first token ($5, else, s/…) is then blocked as a command.
    segments = _split_on_unquoted_ops(command)
    for seg in segments:
        seg = seg.strip()
        if not seg:
            continue
        # Strip redirections so "grep foo file > out.txt" validates as "grep"
        seg_clean = _re.sub(r'2?>+\s*\S*', '', seg).strip()
        seg_clean = _re.sub(r'<\s*\S*', '', seg_clean).strip()
        if not seg_clean:
            continue
        try:
            parts = shlex.split(seg_clean)
        except ValueError:
            # Malformed quoting inside the segment — be permissive; it's
            # probably a bash -c subscript that the shell will handle.
            parts = seg_clean.split()
        if not parts:
            continue
        first_word = parts[0]
        base = _base_cmd(first_word)
        # Redirect artifacts: the operator split on '&' turns '2>&1' into a
        # segment '1', and '&>file'/'>&2' into '>file'/'>2'. These are shell
        # redirections, not executables — skip them (fixes the "Command not
        # allowed: 1" false block that stops any command using '2>&1').
        if base.isdigit() or first_word[0] in '<>':
            continue
        # Variable / field reference ($5, $HOME, ${x}) or brace fragment — not an
        # executable (leftover from an awk/sed program or a $VAR at segment start).
        if first_word[0] in '$}{':
            continue
        # Variable assignment (NAME=value or NAME=$(...)): always allowed
        if _re.match(r'^[A-Za-z_][A-Za-z0-9_]*=', first_word):
            continue
        # Shell builtins: always allowed (part of the shell itself)
        if base in _SHELL_BUILTINS:
            continue
        if base not in _ALLOWED_COMMANDS_LOWER:
            blocked.append(base)
    return len(blocked) == 0, blocked


def is_command_allowed(command: str) -> bool:
    """Check that every executable in a pipeline is in ALLOWED_COMMANDS.

    For shell interpreters (bash, sh, zsh) called with -c, validates the
    -c argument content recursively so that bash -c "cmd1 | cmd2 && cmd3"
    works correctly without the outer re.split misinterpreting operators
    inside shell-quoted strings.

    Returns False if any pipeline segment uses a disallowed command.
    """
    if not command or not command.strip():
        return False

    # Parse the outer command to get proper token list (respects quoting)
    try:
        outer_parts = shlex.split(command)
    except ValueError:
        return False
    if not outer_parts:
        return False

    base = _base_cmd(outer_parts[0])

    # Shell interpreter with -c: validate the -c content recursively.
    # This avoids splitting on operators inside the quoted -c argument.
    if base in _SHELL_INTERPRETERS and '-c' in outer_parts:
        c_idx = outer_parts.index('-c')
        if c_idx + 1 < len(outer_parts):
            script_content = outer_parts[c_idx + 1]
            allowed, _ = _validate_segments(script_content)
            return allowed
        # bash -c with no following argument: just validate the interpreter
        return base in _ALLOWED_COMMANDS_LOWER

    # Normal command or pipeline: split on operators and validate each segment.
    # re.split is safe here because we're not inside a bash -c quoted string.
    allowed, _ = _validate_segments(command)
    return allowed


def _get_blocked_commands(command: str) -> list:
    """Return list of disallowed command names — used for error messages."""
    try:
        outer_parts = shlex.split(command)
    except ValueError:
        return [command.split()[0]] if command.split() else []
    if not outer_parts:
        return []
    base = _base_cmd(outer_parts[0])
    if base in _SHELL_INTERPRETERS and '-c' in outer_parts:
        c_idx = outer_parts.index('-c')
        if c_idx + 1 < len(outer_parts):
            _, blocked = _validate_segments(outer_parts[c_idx + 1])
            return blocked
    _, blocked = _validate_segments(command)
    return blocked

# =============================================================================
# SECTION 1: FILESYSTEM OPERATIONS
# =============================================================================

@mcp.tool()
async def list_directory(path: str = ".", show_hidden: bool = False, details: bool = False) -> str:
    """List contents of a directory.

    Args:
        path: Directory path to list (default: current directory)
        show_hidden: Include hidden files/directories starting with '.' (default: False)
        details: Show detailed file information (default: False)

    Returns:
        JSON with directory contents
    """
    try:
        dir_path = Path(path).expanduser().resolve()

        if not dir_path.exists():
            return format_result({"error": f"Directory does not exist: {dir_path}"})

        if not dir_path.is_dir():
            return format_result({"error": f"Path is not a directory: {path}"})

        def human_size(n):
            for unit in ('B', 'K', 'M', 'G'):
                if n < 1024:
                    return f"{n:.0f}{unit}" if unit == 'B' else f"{n:.1f}{unit}"
                n /= 1024
            return f"{n:.1f}T"

        # ANSI color codes for terminal-style output
        BOLD = "\033[1m"
        BLUE = "\033[1;34m"
        CYAN = "\033[1;36m"
        GREEN = "\033[1;32m"
        YELLOW = "\033[33m"
        MAGENTA = "\033[35m"
        DIM = "\033[2m"
        RESET = "\033[0m"

        ext_colors = {
            '.py': GREEN, '.sh': GREEN, '.js': GREEN, '.ts': GREEN, '.tsx': GREEN,
            '.csv': YELLOW, '.dat': YELLOW, '.xy': YELLOW, '.txt': YELLOW,
            '.tif': MAGENTA, '.tiff': MAGENTA, '.ge': MAGENTA, '.ge2': MAGENTA,
            '.ge3': MAGENTA, '.ge4': MAGENTA, '.ge5': MAGENTA, '.h5': MAGENTA,
            '.hdf': MAGENTA, '.zarr': MAGENTA,
            '.md': DIM, '.log': DIM, '.lock': DIM,
            '.json': CYAN, '.toml': CYAN, '.yaml': CYAN, '.yml': CYAN,
            '.pdf': "\033[31m",
        }

        dir_names = []
        file_entries = []   # (name, colored_name) for column layout
        symlink_lines = []
        hidden_count = 0
        try:
            terminal_width = os.get_terminal_size().columns
        except OSError:
            terminal_width = 120

        for item in sorted(dir_path.iterdir(), key=lambda p: p.name.lower()):
            if item.name.startswith('.'):
                if not show_hidden:
                    hidden_count += 1
                    continue
            if item.is_symlink():
                try:
                    target = os.readlink(item)
                    resolved = item.resolve()
                    if resolved.is_dir():
                        symlink_lines.append(f"{CYAN}{item.name}/{RESET} {DIM}-> {target}{RESET}")
                    else:
                        symlink_lines.append(f"{item.name} {DIM}-> {target}{RESET}")
                except OSError:
                    symlink_lines.append(f"{item.name} {DIM}[broken]{RESET}")
            elif item.is_dir():
                dir_names.append(item.name + "/")
            else:
                ext = item.suffix.lower()
                color = ext_colors.get(ext, "")
                reset = RESET if color else ""
                file_entries.append((item.name, f"{color}{item.name}{reset}"))

        def _render_columns(items, colored_items, indent=2):
            """Render items in multi-column layout like ls."""
            if not items:
                return ""
            max_len = max(len(name) for name in items) + 2
            cols = max(1, (terminal_width - indent) // max_len)
            rows = []
            for i in range(0, len(items), cols):
                row_items = []
                for j in range(i, min(i + cols, len(items))):
                    plain_name = items[j]
                    colored = colored_items[j]
                    padding = max_len - len(plain_name)
                    row_items.append(colored + " " * padding)
                rows.append(" " * indent + "".join(row_items).rstrip())
            return "\n".join(rows)

        output = f"{BOLD}{str(dir_path.absolute())}{RESET}\n"

        if dir_names:
            colored_dirs = [f"{BLUE}{BOLD}{name}{RESET}" for name in dir_names]
            output += _render_columns(dir_names, colored_dirs) + "\n"

        if symlink_lines:
            for sl in symlink_lines:
                output += f"  {sl}\n"

        if file_entries:
            if dir_names or symlink_lines:
                output += "\n"
            plain_names = [e[0] for e in file_entries]
            colored_names = [e[1] for e in file_entries]
            output += _render_columns(plain_names, colored_names)

        if hidden_count:
            output += f"\n  {DIM}({hidden_count} hidden items — use show_hidden=True){RESET}"

        total = len(dir_names) + len(file_entries) + len(symlink_lines)
        output += f"\n  {DIM}{len(dir_names)} directories, {len(file_entries)} files{RESET}"

        return format_result({
            "tool": "list_directory",
            "path": str(dir_path.absolute()),
            "listing": output,
            "dirs": dir_names,
            "files": [e[0] for e in file_entries],
        })

    except Exception as e:
        return format_result({"error": f"Error listing directory: {str(e)}"})

@mcp.tool()
async def read_file(file_path: str, encoding: str = "utf-8", max_size: int = 1024000) -> str:
    """Read contents of a text file.

    Args:
        file_path: Path to the file to read
        encoding: Text encoding (default: utf-8)
        max_size: Maximum file size to read in bytes (default: 1MB)

    Returns:
        JSON with file contents
    """
    try:
        path = Path(file_path).expanduser().resolve()

        if not path.exists():
            return format_result({"error": f"File does not exist: {path}"})

        if not path.is_file():
            return format_result({"error": f"Path is not a file: {file_path}"})

        # Transparently route document formats to read_document so the caller
        # gets text even if it reached for the wrong tool (PDFs/Office files are
        # binary and would otherwise fail the binary check below).
        _DOC_EXTS = {".pdf", ".pptx", ".docx", ".xlsx", ".odt", ".odp",
                     ".ods", ".rtf"}
        if path.suffix.lower() in _DOC_EXTS:
            return await read_document(str(path))

        file_size = path.stat().st_size
        if file_size > max_size:
            return format_result({"error": f"File too large ({file_size} bytes > {max_size} bytes limit)"})

        # Detect if file is binary
        with open(path, 'rb') as f:
            sample = f.read(1024)
            if b'\x00' in sample:
                return format_result({
                    "error": f"File appears to be binary: {file_path}. "
                             "If this is a document, use read_document instead."})

        # Read text file
        with open(path, 'r', encoding=encoding) as f:
            content = f.read()

        return format_result({
            "tool": "read_file",
            "file_path": str(path.absolute()),
            "size": file_size,
            "encoding": encoding,
            "line_count": content.count('\n') + 1,
            "content": content[:10000],
            "truncated": len(content) > 10000
        })

    except Exception as e:
        return format_result({"error": f"Error reading file: {str(e)}"})


# ── read_document helpers ───────────────────────────────────────────────────
# Each returns (text, meta_dict, error_str). Prefer a real parser library when
# installed; fall back to a dependency-free zipfile+XML scrape for OOXML/ODF so
# the tool works on a fresh checkout before `uv sync` adds the nicer libs.

def _ooxml_xml_fallback(path, member_prefix: str):
    """Universal OOXML/ODF text scrape: strip tags from matching zip members.

    Lossy (no layout) but dependency-free — used when python-pptx/docx/openpyxl
    are not installed. member_prefix selects the relevant parts of the archive
    (e.g. 'ppt/slides/', 'word/document.xml', 'content.xml').
    """
    import zipfile
    import re as _re
    try:
        chunks = []
        with zipfile.ZipFile(str(path)) as z:
            names = sorted(
                n for n in z.namelist()
                if n.startswith(member_prefix) and n.endswith(".xml")
            )
            # ODF stores everything in a single content.xml (exact match).
            if not names and member_prefix.endswith(".xml"):
                names = [member_prefix] if member_prefix in z.namelist() else []
            for n in names:
                raw = z.read(n).decode("utf-8", errors="ignore")
                # Turn paragraph/row/break tags into whitespace, then drop tags.
                raw = _re.sub(r"</(w:p|a:p|text:p|row|tr)>", "\n", raw)
                raw = _re.sub(r"<[^>]+>", " ", raw)
                raw = _re.sub(r"[ \t]+", " ", raw)
                raw = _re.sub(r"\n\s+", "\n", raw).strip()
                if raw:
                    chunks.append(raw)
        return "\n".join(chunks), {}, None
    except Exception as e:
        return None, {}, f"zip/xml fallback failed: {e}"


def _extract_pdf(path):
    try:
        try:
            from pypdf import PdfReader  # maintained successor
        except ImportError:
            from PyPDF2 import PdfReader  # deprecated but widely installed
    except ImportError:
        return None, {}, "PDF support needs pypdf or PyPDF2 (pip install pypdf)"
    try:
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            try:
                pages.append(page.extract_text() or "")
            except Exception:
                pages.append("")
        return "\n\n".join(pages), {"pages": len(reader.pages)}, None
    except Exception as e:
        return None, {}, f"PDF parse failed: {e}"


def _extract_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        return _ooxml_xml_fallback(path, "ppt/slides/")
    try:
        prs = Presentation(str(path))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in para.runs)
                        if line.strip():
                            out.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        out.append("\t".join(c.text for c in row.cells))
            notes = getattr(slide, "notes_slide", None) if slide.has_notes_slide else None
            if notes and notes.notes_text_frame and notes.notes_text_frame.text.strip():
                out.append(f"[notes] {notes.notes_text_frame.text.strip()}")
        return "\n".join(out), {"slides": len(prs.slides)}, None
    except Exception:
        return _ooxml_xml_fallback(path, "ppt/slides/")


def _extract_docx(path):
    try:
        from docx import Document
    except ImportError:
        return _ooxml_xml_fallback(path, "word/document.xml")
    try:
        doc = Document(str(path))
        out = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                out.append("\t".join(c.text for c in row.cells))
        return "\n".join(out), {"paragraphs": len(doc.paragraphs)}, None
    except Exception:
        return _ooxml_xml_fallback(path, "word/document.xml")


def _extract_xlsx(path):
    try:
        import openpyxl
    except ImportError:
        return _ooxml_xml_fallback(path, "xl/")
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"--- Sheet: {ws.title} ---")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    out.append("\t".join(cells))
        return "\n".join(out), {"sheets": len(wb.worksheets)}, None
    except Exception:
        return _ooxml_xml_fallback(path, "xl/")


def _extract_html(path):
    import re as _re
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    raw = _re.sub(r"(?is)<(script|style).*?</\1>", " ", raw)
    raw = _re.sub(r"(?i)</(p|div|br|li|tr|h[1-6])>", "\n", raw)
    raw = _re.sub(r"<[^>]+>", " ", raw)
    raw = _re.sub(r"[ \t]+", " ", raw)
    return _re.sub(r"\n\s+", "\n", raw).strip(), {}, None


def _extract_rtf(path):
    import re as _re
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    raw = _re.sub(r"\\par[d]?", "\n", raw)
    raw = _re.sub(r"\\'[0-9a-fA-F]{2}", "", raw)
    raw = _re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw)
    raw = raw.replace("{", "").replace("}", "")
    return _re.sub(r"\n\s+", "\n", raw).strip(), {}, None


@mcp.tool()
async def read_document(file_path: str, max_chars: int = 100000,
                        max_size: int = 52428800) -> str:
    """Read text from documents of many formats — logbooks, slides, spreadsheets.

    Use this (not read_file) for binary/office documents: experiment logbooks
    (PDF), methodology slides (PPTX), reports (DOCX), data tables (XLSX),
    OpenDocument files, HTML, and RTF. read_file only handles plain text.

    Supported: .pdf, .pptx, .docx, .xlsx, .odt/.odp/.ods, .html/.htm, .rtf,
    .csv/.tsv, and any plain-text/code/markup file. Office/ODF formats use a
    proper parser when installed (pypdf, python-pptx, python-docx, openpyxl)
    and fall back to a dependency-free zip+XML text scrape otherwise.

    Args:
        file_path: Path to the document.
        max_chars: Max characters of extracted text to return (default 100000).
        max_size:  Max file size to open, bytes (default 50 MB).

    Returns:
        JSON with extracted text, detected format, and any format metadata
        (page/slide/sheet counts), plus a truncated flag.
    """
    try:
        path = Path(file_path).expanduser().resolve()
        if not path.exists():
            return format_result({"error": f"File does not exist: {path}"})
        if not path.is_file():
            return format_result({"error": f"Path is not a file: {file_path}"})

        size = path.stat().st_size
        if size > max_size:
            return format_result({"error": f"File too large ({size} bytes > {max_size} limit)"})

        ext = path.suffix.lower()
        TEXT_EXT = {".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json",
                    ".yaml", ".yml", ".xml", ".log", ".ini", ".cfg", ".conf",
                    ".toml", ".py", ".c", ".h", ".cpp", ".cc", ".hpp", ".sh",
                    ".js", ".ts", ".tex", ".bib", ".m", ".f", ".f90"}

        if ext == ".pdf":
            text, meta, err = _extract_pdf(path)
        elif ext == ".pptx":
            text, meta, err = _extract_pptx(path)
        elif ext == ".docx":
            text, meta, err = _extract_docx(path)
        elif ext == ".xlsx":
            text, meta, err = _extract_xlsx(path)
        elif ext in (".odt", ".odp", ".ods"):
            text, meta, err = _ooxml_xml_fallback(path, "content.xml")
        elif ext in (".html", ".htm"):
            text, meta, err = _extract_html(path)
        elif ext == ".rtf":
            text, meta, err = _extract_rtf(path)
        elif ext in TEXT_EXT:
            text, meta, err = path.read_text(encoding="utf-8", errors="replace"), {}, None
        else:
            # Unknown extension: try as text, but refuse true binary.
            with open(path, "rb") as f:
                if b"\x00" in f.read(1024):
                    return format_result({
                        "error": f"Unsupported binary format '{ext}'. Supported: "
                                 "pdf, pptx, docx, xlsx, odt/odp/ods, html, rtf, "
                                 "csv/tsv, and plain text/code.",
                        "file_path": str(path),
                    })
            text, meta, err = path.read_text(encoding="utf-8", errors="replace"), {}, None

        if err:
            return format_result({"tool": "read_document", "status": "error",
                                  "file_path": str(path), "format": ext, "error": err})

        text = text or ""
        return format_result({
            "tool": "read_document",
            "file_path": str(path),
            "format": ext.lstrip(".") or "text",
            "size": size,
            "char_count": len(text),
            "truncated": len(text) > max_chars,
            **meta,
            "content": text[:max_chars],
        })

    except Exception as e:
        return format_result({"error": f"Error reading document: {str(e)}"})


@mcp.tool()
async def write_file(file_path: str, content: str, encoding: str = "utf-8", append: bool = False) -> str:
    """Write content to a file.

    Args:
        file_path: Path to the file to write
        content: Content to write
        encoding: Text encoding (default: utf-8)
        append: Append to file instead of overwriting (default: False)

    Returns:
        JSON with operation status
    """
    try:
        path = Path(file_path).expanduser().resolve()
        mode = 'a' if append else 'w'

        with open(path, mode, encoding=encoding) as f:
            f.write(content)

        return format_result({
            "tool": "write_file",
            "status": "success",
            "file_path": str(path.absolute()),
            "bytes_written": len(content.encode(encoding)),
            "mode": "append" if append else "overwrite"
        })

    except Exception as e:
        return format_result({"error": f"Error writing file: {str(e)}"})

@mcp.tool()
async def get_file_info(file_path: str) -> str:
    """Get detailed information about a file or directory.

    Args:
        file_path: Path to the file or directory

    Returns:
        JSON with detailed file information
    """
    try:
        path = Path(file_path).expanduser().resolve()

        if not path.exists():
            return format_result({"error": f"Path does not exist: {path}"})

        info = format_file_info(path)
        info["tool"] = "get_file_info"

        return format_result(info)

    except Exception as e:
        return format_result({"error": f"Error getting file info: {str(e)}"})

# =============================================================================
# SECTION 2: COMMAND EXECUTION
# =============================================================================

@mcp.tool()
async def run_command(command: str, working_dir: str = None, timeout: int = 120) -> str:
    """Execute a shell command with full bash capabilities.

    Supports pipes (|), redirections (>, >>), &&, ||, semicolons, subshells,
    and bash -c "..." multi-command scripts — equivalent to Claude Code's
    Bash tool for beamline analysis tasks. Every executable in a pipeline
    is validated against an allowed-command list before execution.

    Use freely for: grep, awk, sed, find, wc, sort, uniq, diff, head, tail,
    cat, du, stat, ls, python3 scripts, MIDAS executables, and pipelines
    combining any of the above.

    Args:
        command: Shell command to execute. Full bash syntax supported.
                 Examples:
                   "find /path -name '*.h5' | wc -l"
                   "grep -n 'Wavelength' params.txt | head -5"
                   "awk -F, 'NR>1{print $3}' Grains.csv | sort | uniq -c"
                   "bash -c 'for f in /path/*.h5; do echo $f; done'"
        working_dir: Working directory (defaults to CWD if not set)
        timeout: Max execution time in seconds (default: 120)

    Returns:
        JSON with stdout, stderr, return_code, success flag, working_dir
    """
    try:
        if not is_command_allowed(command):
            blocked = _get_blocked_commands(command)
            return format_result({
                "error": f"Command not allowed: {', '.join(blocked) if blocked else command.split()[0]}",
                "detail": "Every executable in the pipeline must be in the allowed list.",
                "allowed_commands": "Use check_environment to see allowed commands"
            })

        cwd = working_dir if working_dir and Path(working_dir).exists() else None

        # shell=True enables pipes (|), semicolons (;), &&, ||, redirections
        # (>, >>), and other shell operators.  Security is preserved by
        # is_command_allowed() which validates every pipeline segment's base
        # executable against ALLOWED_COMMANDS before reaching this point.
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=cwd
        )

        return format_result({
            "tool": "run_command",
            "command": command,
            "return_code": result.returncode,
            "stdout": result.stdout,
            "stderr": result.stderr,
            "success": result.returncode == 0,
            "working_dir": str(cwd) if cwd else str(Path.cwd())
        })

    except subprocess.TimeoutExpired:
        return format_result({"error": f"Command timed out after {timeout} seconds"})
    except Exception as e:
        return format_result({"error": f"Error executing command: {str(e)}"})

@mcp.tool()
async def run_remote_command(command: str, host: str = "", remote_dir: str = None,
                             timeout: int = 600) -> str:
    """Execute a shell command on a REMOTE analysis host over SSH.

    Like run_command, but the command runs on another machine — for running
    analysis where the beamline data already lives (e.g. `copland`), instead of
    copying the data to this machine. The command is executed through a remote
    LOGIN shell (`bash -lc`), so the remote user's PATH/profile is sourced and
    MIDAS executables resolve exactly as they would in an interactive session
    there.

    SSH MUST BE KEY-BASED. This runs non-interactively (BatchMode=yes); a
    password prompt would hang, so instead it fails fast (return_code 255) with
    an actionable error. Set up once with `ssh-copy-id <user>@<host>`.

    The same allowlist as run_command is applied to `command` before it is sent,
    so bash/python/MIDAS binaries/ls/rsync/grep/… are permitted and unknown
    executables are refused — the guardrail travels with the command to the
    remote host.

    Args:
        command: Shell command to run on the remote host. Full bash syntax
                 supported (pipes, &&, redirects). Examples:
                   "ls -la /gdata/dm/1ID/2026/pokharel_jul26/data/ge5"
                   "cd /gdata/.../pokharel_jul26 && ff_MIDAS.py -paramFile ff.txt"
                   "nproc && df -h /gdata"
        host: Remote host, e.g. "copland" or "user@copland". Defaults to
              $APEXA_ANALYSIS_HOST, else "copland".
        remote_dir: Directory to `cd` into on the remote host before running.
        timeout: Max seconds (default 600; remote analysis can be long).

    Returns:
        JSON with stdout, stderr, return_code, success, host, command.
    """
    try:
        target = host or os.environ.get("APEXA_ANALYSIS_HOST", "copland")

        if not is_command_allowed(command):
            blocked = _get_blocked_commands(command)
            return format_result({
                "error": f"Command not allowed: {', '.join(blocked) if blocked else command.split()[0]}",
                "detail": "Every executable in the remote command must be in the allowed list.",
                "allowed_commands": "Use check_environment to see allowed commands"
            })

        # Run over SSH via the shared transport (login shell `bash -lc`, BatchMode,
        # ConnectTimeout, rc==255 = ssh-layer failure). Factored into
        # apexa_remote_exec so the midas server's per-tool routing shares one
        # implementation; behavior and response keys here are unchanged.
        res = remote_run(target, command, remote_dir=remote_dir, timeout=timeout)

        if res.get("timed_out"):
            return format_result({
                "error": f"Remote command timed out after {timeout} seconds on '{target}'"
            })

        # ssh exit 255 = ssh-layer failure (auth/connectivity), not the remote
        # command's own exit code. Classify it so the agent gets an actionable
        # message instead of a bare non-zero.
        if res.get("ssh_failed"):
            return format_result({
                "tool": "run_remote_command",
                "host": target,
                "command": command,
                "return_code": 255,
                "stderr": res.get("stderr", ""),
                "success": False,
                "error": ssh_hint(target),
            })

        return format_result({
            "tool": "run_remote_command",
            "host": target,
            "command": command,
            "remote_dir": remote_dir,
            "return_code": res.get("return_code"),
            "stdout": res.get("stdout", ""),
            "stderr": res.get("stderr", ""),
            "success": res.get("success", False),
        })

    except Exception as e:
        return format_result({
            "error": f"Error executing remote command: {str(e)}"
        })

@mcp.tool()
async def check_environment() -> str:
    """Check system environment and available tools.

    Returns:
        JSON with system information and available MIDAS tools
    """
    try:
        info = {
            "tool": "check_environment",
            "system": {
                "platform": platform.system(),
                "python_version": platform.python_version(),
                "architecture": platform.machine(),
                "hostname": platform.node()
            },
            "paths": {
                "current_dir": str(Path.cwd()),
                "home_dir": str(Path.home()),
                "python_executable": sys.executable
            },
            "available_commands": sorted(list(ALLOWED_COMMANDS)),
            "optional_libraries": {
                "scipy": SCIPY_AVAILABLE,
                "fabio": SCIPY_AVAILABLE
            }
        }

        # Check for MIDAS installation (same search as midas server)
        # expanduser() so a literal "~" in MIDAS_PATH (e.g. MIDAS_PATH=~/opt/MIDAS_canonical
        # in .env) resolves — Path("~/...").exists() is always False without it.
        midas_path = os.environ.get("MIDAS_PATH")
        if midas_path:
            midas_path = str(Path(midas_path).expanduser())
        if midas_path and Path(midas_path).exists():
            info["midas"] = {
                "installed": True,
                "path": midas_path
            }
        else:
            # Search common locations (mirrors find_midas_installation in midas server)
            found = None
            for subdir in ["Git", "git", "src", "opt", ""]:
                for name in ["MIDAS", "midas"]:
                    candidate = Path.home() / subdir / name if subdir else Path.home() / name
                    autocal = candidate / "utils" / "AutoCalibrateZarr.py"
                    if candidate.exists() and autocal.exists():
                        found = str(candidate)
                        break
                if found:
                    break
            if found:
                info["midas"] = {
                    "installed": True,
                    "path": found
                }
            else:
                info["midas"] = {
                    "installed": False,
                    "message": "Set MIDAS_PATH environment variable"
                }

        return format_result(info)

    except Exception as e:
        return format_result({"error": f"Error checking environment: {str(e)}"})

# =============================================================================
# SECTION 3: X-RAY UTILITIES & CALCULATIONS (using xrayutilities)
# =============================================================================

# Canonical Planck constant × c for E[keV] <-> lambda[Å]:  lambda = HC / E.
# This exact value is the beamline's own energy2lambda script constant
# (AC, cchuang@anl.gov, 2025-04-10) and agrees with xrayutilities.en2lam to
# 7 significant figures. Pinned here as the SINGLE source of truth so any
# wavelength APEXA derives (e.g. the calibration --wavelength) equals the
# beamline's other tools exactly, rather than a model-computed approximation.
HC_KEV_ANGSTROM = 12.398419057638671

# Element K-edge energies (keV) for edge-tuned runs. Source: the same beamline
# energy2lambda script. Used ONLY when a run is deliberately tuned to an element's
# absorption edge — a *sample* made of gold does NOT imply the Au edge; the beam
# energy comes from the experiment (e.g. "96keV" in the filename), not the element.
ELEMENT_EDGE_KEV = {
    "ho": 55.615, "yb": 61.332, "lu": 63.314, "hf": 65.350,
    "ta": 67.411, "w": 69.525, "re": 71.676, "ir": 76.112,
    "au": 80.726, "pb": 88.005, "bi": 90.529,
}


def _lambda_from_energy_kev(energy_kev: float) -> float:
    """Wavelength in Å from energy in keV using the canonical beamline constant."""
    return HC_KEV_ANGSTROM / energy_kev


def _energy_kev_from_lambda(wavelength_angstroms: float) -> float:
    """Energy in keV from wavelength in Å using the canonical beamline constant."""
    return HC_KEV_ANGSTROM / wavelength_angstroms


@mcp.tool()
async def xray_calculate(
    calculation_type: str,
    h: int = None,
    k: int = None,
    l: int = None,
    material: str = None,
    lattice_a: float = None,
    two_theta_degrees: float = None,
    wavelength_angstroms: float = None,
    energy_kev: float = None,
    element: str = None,
    d_spacing: float = None,
    measured_d: float = None,
    reference_d: float = None
) -> str:
    """Universal X-ray calculation tool using xrayutilities library.

    🎯 CRITICAL: This tool MUST be used for ALL X-ray calculations.
    Uses domain-specific xrayutilities package for accurate scientific calculations.

    Supported calculation types:
    - "d_from_hkl": Calculate d-spacing from Miller indices
    - "d_from_angle": Calculate d-spacing from 2θ (Bragg's law)
    - "angle_from_d": Calculate 2θ from d-spacing
    - "energy_to_wavelength": Convert energy (keV) to wavelength (Å). Accepts
        energy_kev directly, OR element=<symbol> to use that element's K-edge
        energy (edge-tuned runs only — see list_element_edges).
    - "wavelength_to_energy": Convert wavelength (Å) to energy (keV)
    - "list_element_edges": List the element K-edge energies (keV) known to APEXA
    - "strain": Calculate strain from measured and reference d-spacings
    - "list_materials": List available materials in xrayutilities

    NOTE: energy<->wavelength uses the beamline's canonical constant
    (lambda = 12.398419057638671 / E[keV]). ALWAYS derive a calibration/FF
    wavelength through this tool — do not hand-compute it.

    Args:
        calculation_type: Type of calculation (see above)
        h, k, l: Miller indices (for d_from_hkl)
        material: Material name (e.g., "Fe", "Si", "Al") from xrayutilities
        lattice_a: Lattice parameter in Å (for custom materials)
        two_theta_degrees: Diffraction angle 2θ in degrees
        wavelength_angstroms: X-ray wavelength in Å
        energy_kev: X-ray energy in keV
        d_spacing: d-spacing in Å
        measured_d: Measured d-spacing in Å (for strain)
        reference_d: Reference d-spacing in Å (for strain)

    Returns:
        JSON with calculation results

    Examples:
        xray_calculate("d_from_hkl", h=1, k=1, l=0, material="Fe")
        xray_calculate("d_from_angle", two_theta_degrees=12.5, wavelength_angstroms=0.202)
        xray_calculate("energy_to_wavelength", energy_kev=61.332)
    """
    try:
        if not XRAYUTILITIES_AVAILABLE:
            return format_result({
                "error": "xrayutilities library not available",
                "install": "pip install xrayutilities",
                "fallback": "Some basic calculations available without xrayutilities"
            })

        # D-SPACING FROM MILLER INDICES
        if calculation_type == "d_from_hkl":
            if h is None or k is None or l is None:
                return format_result({"error": "h, k, l Miller indices required"})

            if material:
                # Use xrayutilities materials database
                try:
                    mat = getattr(xumaterials, material)
                    d = mat.planeDistance(h, k, l)

                    return format_result({
                        "tool": "xray_calculate",
                        "library": "xrayutilities",
                        "calculation": "d_from_hkl",
                        "inputs": {
                            "miller_indices": f"({h}{k}{l})",
                            "material": material,
                            "lattice_parameters": {
                                "a": mat.a,
                                "b": mat.b if hasattr(mat, 'b') else mat.a,
                                "c": mat.c if hasattr(mat, 'c') else mat.a
                            }
                        },
                        "result": {
                            "d_spacing_angstroms": round(d, 6),
                            "d_spacing_nm": round(d / 10, 6)
                        }
                    })
                except AttributeError:
                    return format_result({
                        "error": f"Material '{material}' not found in xrayutilities database",
                        "suggestion": "Use calculation_type='list_materials' to see available materials"
                    })

            elif lattice_a:
                # Calculate for cubic system with custom lattice
                d = lattice_a / np.sqrt(h**2 + k**2 + l**2)
                return format_result({
                    "tool": "xray_calculate",
                    "calculation": "d_from_hkl",
                    "inputs": {"h": h, "k": k, "l": l, "lattice_a": lattice_a},
                    "result": {"d_spacing_angstroms": round(d, 6)},
                    "note": "Cubic system assumed"
                })
            else:
                return format_result({"error": "Provide either 'material' or 'lattice_a'"})

        # D-SPACING FROM ANGLE (Bragg's law)
        elif calculation_type == "d_from_angle":
            if two_theta_degrees is None or wavelength_angstroms is None:
                return format_result({"error": "two_theta_degrees and wavelength_angstroms required"})

            theta_rad = np.radians(two_theta_degrees / 2.0)
            d = wavelength_angstroms / (2.0 * np.sin(theta_rad))

            return format_result({
                "tool": "xray_calculate",
                "library": "xrayutilities",
                "calculation": "d_from_angle",
                "formula": "d = λ / (2 * sin(θ))",
                "inputs": {
                    "two_theta_degrees": two_theta_degrees,
                    "wavelength_angstroms": wavelength_angstroms
                },
                "result": {"d_spacing_angstroms": round(d, 6)}
            })

        # ANGLE FROM D-SPACING
        elif calculation_type == "angle_from_d":
            if d_spacing is None or wavelength_angstroms is None:
                return format_result({"error": "d_spacing and wavelength_angstroms required"})

            sin_theta = wavelength_angstroms / (2.0 * d_spacing)
            if sin_theta > 1.0:
                return format_result({
                    "error": "No diffraction possible",
                    "reason": f"λ/(2d) = {sin_theta:.3f} > 1.0"
                })

            two_theta = 2.0 * np.degrees(np.arcsin(sin_theta))

            return format_result({
                "tool": "xray_calculate",
                "calculation": "angle_from_d",
                "formula": "2θ = 2 * arcsin(λ / (2d))",
                "inputs": {"d_spacing": d_spacing, "wavelength_angstroms": wavelength_angstroms},
                "result": {"two_theta_degrees": round(two_theta, 4)}
            })

        # ENERGY TO WAVELENGTH
        elif calculation_type == "energy_to_wavelength":
            # Resolve energy: explicit energy_kev wins; else an element symbol
            # maps to its K-edge energy (edge-tuned runs only).
            _edge_note = None
            if energy_kev is None and element:
                _el = element.strip().lower()
                if _el not in ELEMENT_EDGE_KEV:
                    return format_result({
                        "error": f"Unknown element '{element}'",
                        "known_elements": sorted(ELEMENT_EDGE_KEV),
                        "hint": "Pass energy_kev directly, or use calculation_type='list_element_edges'.",
                    })
                energy_kev = ELEMENT_EDGE_KEV[_el]
                _edge_note = (f"Using {_el.capitalize()} K-edge = {energy_kev} keV. "
                              "This is an edge-tuned energy — for a normal run pass the "
                              "actual beam energy (e.g. from the filename), not the sample element.")
            elif energy_kev is not None and element:
                _edge_note = "Both energy_kev and element given — using explicit energy_kev."
            if energy_kev is None:
                return format_result({"error": "energy_kev or element required"})

            # Canonical beamline constant (matches xrayutilities.en2lam to 7 figs).
            wavelength = _lambda_from_energy_kev(energy_kev)

            _out = {
                "tool": "xray_calculate",
                "calculation": "energy_to_wavelength",
                "constant": "lambda = 12.398419057638671 / E[keV] (beamline energy2lambda)",
                "inputs": {"energy_kev": energy_kev,
                           **({"element": element} if element else {})},
                "result": {
                    "wavelength_angstroms": round(wavelength, 6),
                    "energy_eV": energy_kev * 1000,
                },
            }
            if _edge_note:
                _out["note"] = _edge_note
            return format_result(_out)

        # WAVELENGTH TO ENERGY
        elif calculation_type == "wavelength_to_energy":
            if wavelength_angstroms is None:
                return format_result({"error": "wavelength_angstroms required"})

            # Canonical beamline constant (matches xrayutilities.lam2en to 7 figs).
            energy_kev = _energy_kev_from_lambda(wavelength_angstroms)

            return format_result({
                "tool": "xray_calculate",
                "calculation": "wavelength_to_energy",
                "constant": "E[keV] = 12.398419057638671 / lambda[Å] (beamline energy2lambda)",
                "inputs": {"wavelength_angstroms": wavelength_angstroms},
                "result": {
                    "energy_kev": round(energy_kev, 6),
                    "energy_eV": round(energy_kev * 1000, 2)
                }
            })

        # LIST ELEMENT K-EDGE ENERGIES
        elif calculation_type == "list_element_edges":
            return format_result({
                "tool": "xray_calculate",
                "calculation": "list_element_edges",
                "source": "beamline energy2lambda script (AC, cchuang@anl.gov)",
                "element_k_edges_keV": ELEMENT_EDGE_KEV,
                "note": "Edge-tuned energies only. A gold sample measured at 96 keV "
                        "uses 96 keV, NOT the Au edge (80.726 keV).",
            })

        # STRAIN CALCULATION
        elif calculation_type == "strain":
            if measured_d is None or reference_d is None:
                return format_result({"error": "measured_d and reference_d required"})

            strain = (measured_d - reference_d) / reference_d
            microstrain = strain * 1e6

            return format_result({
                "tool": "xray_calculate",
                "calculation": "strain",
                "formula": "ε = (d_measured - d_ref) / d_ref",
                "inputs": {"measured_d": measured_d, "reference_d": reference_d},
                "result": {
                    "strain": round(strain, 9),
                    "microstrain": round(microstrain, 2),
                    "percent": round(strain * 100, 4),
                    "type": "tension" if strain > 0 else "compression" if strain < 0 else "no strain"
                }
            })

        # LIST AVAILABLE MATERIALS
        elif calculation_type == "list_materials":
            # Get all materials from xrayutilities
            materials = [name for name in dir(xumaterials)
                        if not name.startswith('_') and name[0].isupper()]

            return format_result({
                "tool": "xray_calculate",
                "library": "xrayutilities.materials",
                "calculation": "list_materials",
                "available_materials": sorted(materials[:50]),  # First 50
                "total_count": len(materials),
                "note": "Use material name in d_from_hkl calculation"
            })

        else:
            return format_result({
                "error": f"Unknown calculation_type: {calculation_type}",
                "valid_types": ["d_from_hkl", "d_from_angle", "angle_from_d",
                               "energy_to_wavelength", "wavelength_to_energy",
                               "list_element_edges", "strain", "list_materials"]
            })

    except Exception as e:
        return format_result({"error": f"X-ray calculation error: {str(e)}"})

@mcp.tool()
async def validate_beamline_parameters(
    energy_kev: float,
    detector_distance_mm: float,
    beam_center_x: float,
    beam_center_y: float,
    pixel_size_um: float = None
) -> str:
    """Validate that beamline experimental parameters are physically reasonable.

    Checks:
    - Energy is in typical synchrotron range
    - Detector distance is reasonable
    - Beam center is within typical detector bounds
    - Pixel size is reasonable if provided

    Args:
        energy_kev: X-ray energy in keV
        detector_distance_mm: Sample-to-detector distance in millimeters
        beam_center_x: Beam center X coordinate in pixels
        beam_center_y: Beam center Y coordinate in pixels
        pixel_size_um: Detector pixel size in micrometers (optional)

    Returns:
        JSON with validation results and any warnings
    """
    try:
        issues = []
        warnings = []

        # Energy validation (typical range: 5-150 keV for synchrotrons)
        if energy_kev < 5 or energy_kev > 150:
            issues.append(f"Energy {energy_kev} keV is outside typical synchrotron range (5-150 keV)")
        elif energy_kev < 10 or energy_kev > 120:
            warnings.append(f"Energy {energy_kev} keV is unusual (typical: 10-120 keV)")

        # Detector distance (typical: 50-2000 mm for HEDM)
        if detector_distance_mm < 50 or detector_distance_mm > 2000:
            issues.append(f"Detector distance {detector_distance_mm} mm seems unusual (typical: 50-2000 mm)")

        # Beam center (typical detector: 2048x2048 pixels, but can be 4096x4096)
        if beam_center_x < 0 or beam_center_x > 5000:
            issues.append(f"Beam center X={beam_center_x} outside typical range (0-5000 pixels)")
        if beam_center_y < 0 or beam_center_y > 5000:
            issues.append(f"Beam center Y={beam_center_y} outside typical range (0-5000 pixels)")

        # Pixel size validation (typical: 50-200 μm)
        if pixel_size_um is not None:
            if pixel_size_um < 10 or pixel_size_um > 500:
                warnings.append(f"Pixel size {pixel_size_um} μm is unusual (typical: 50-200 μm)")

        # Calculate wavelength for reference (xrayutilities-backed)
        from apexa_units import kev_to_angstrom
        wavelength = kev_to_angstrom(energy_kev)

        result = {
            "tool": "validate_beamline_parameters",
            "valid": len(issues) == 0,
            "issues": issues,
            "warnings": warnings,
            "parameters": {
                "energy_kev": energy_kev,
                "wavelength_angstroms": round(wavelength, 6),
                "detector_distance_mm": detector_distance_mm,
                "beam_center": [beam_center_x, beam_center_y]
            }
        }

        if pixel_size_um:
            result["parameters"]["pixel_size_um"] = pixel_size_um

        if len(issues) == 0 and len(warnings) == 0:
            result["message"] = "All parameters are within reasonable ranges"

        return format_result(result)

    except Exception as e:
        return format_result({"error": f"Error validating parameters: {str(e)}"})

@mcp.tool()
async def list_common_calibrants() -> str:
    """List commonly used calibration materials with their properties.

    Returns:
        JSON with standard calibrants (CeO2, LaB6, Si, etc.) and their lattice parameters
    """
    calibrants = {
        "CeO2": {
            "name": "Cerium Dioxide (Ceria)",
            "formula": "CeO2",
            "crystal_system": "Cubic",
            "space_group": "Fm-3m (225)",
            "lattice_parameter_a": 5.411,
            "common_use": "Standard calibrant for HEDM, powder diffraction",
            "major_peaks_hkl": ["111", "200", "220", "311", "222", "400", "331", "420"]
        },
        "LaB6": {
            "name": "Lanthanum Hexaboride",
            "formula": "LaB6",
            "crystal_system": "Cubic",
            "space_group": "Pm-3m (221)",
            "lattice_parameter_a": 4.156,
            "common_use": "High-angle calibration, instrumental broadening",
            "major_peaks_hkl": ["100", "110", "111", "200", "210", "211", "220"]
        },
        "Si": {
            "name": "Silicon",
            "formula": "Si",
            "crystal_system": "Cubic",
            "space_group": "Fd-3m (227)",
            "lattice_parameter_a": 5.431,
            "common_use": "NIST standard reference material (SRM 640d)",
            "major_peaks_hkl": ["111", "220", "311", "400", "331", "422", "511", "440"]
        },
        "Al2O3": {
            "name": "Corundum (Alumina)",
            "formula": "Al2O3",
            "crystal_system": "Hexagonal",
            "space_group": "R-3c (167)",
            "lattice_parameter_a": 4.759,
            "lattice_parameter_c": 12.991,
            "common_use": "High-temperature studies",
            "major_peaks_hkl": ["012", "104", "110", "113", "024", "116", "214"]
        },
        "CaF2": {
            "name": "Calcium Fluoride (Fluorite)",
            "formula": "CaF2",
            "crystal_system": "Cubic",
            "space_group": "Fm-3m (225)",
            "lattice_parameter_a": 5.463,
            "common_use": "Calibration standard",
            "major_peaks_hkl": ["111", "220", "311", "400", "331", "422", "511", "440"]
        }
    }

    return format_result({
        "tool": "list_common_calibrants",
        "count": len(calibrants),
        "calibrants": calibrants,
        "note": "Lattice parameters are at room temperature (298K)"
    })

# =============================================================================
# SERVER ENTRY POINT
# =============================================================================

if __name__ == "__main__":
    if not XRAYUTILITIES_AVAILABLE:
        print("⚠ xrayutilities not available — X-ray calculations limited", file=sys.stderr)
    if not SCIPY_AVAILABLE:
        print("⚠ scipy/fabio not available — image analysis limited", file=sys.stderr)
    mcp.run()
